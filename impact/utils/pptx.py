from io import BytesIO

from django.http import HttpResponse
from pptx.oxml.ns import qn


NOMS_DIAPO = [
    "couverture",
    "introduction",
    "sommaire",
    "information-generale",
    "B1-infos-generales",
    "B1-filiales",
    "B1-sites",
    "B1-certifications",
    "B1-certifications-non-pertinent",
    "B2-non-applicable",
    "B2-participation",
    "B2-non-pertinent",
    "B2-pratiques",
    "B2-pratiques-non-pertinent",
    "C1-produits",
    "C1-relations-affaires",
    "C1-relations-affaires-partiellement-non-pertinent",
    "C2-pratiques-1",
    "C2-pratiques-2",
    "C2-pratiques-3",
    "C2-pratiques-4",
    "C2-pratiques-5",
    "C2-pratiques-non-applicable",
    "environnement",
    "B3-consommation",
    "B3-scope3",
    "B4-pollution-air",
    "B4-pollution-air-non-pertinent",
    "B4-pollution-eau",
    "B4-pollution-eau-non-pertinent",
    "B4-pollution-sol",
    "B4-pollution-sol-non-pertinent",
    "B5-biodiversite-non-pertinent",
    "B5-surface",
    "B5-affection-terres-non-pertinent",
    "B5-affection-terres",
    "B6-prelevements",
    "B6-prelevements-partiellement-non-pertinent",
    "B7-economie-circulaire",
    "B7-classification-dechets-1",
    "B7-classification-dechets-2",
    "B7-gestion-dechets",
    "B7-gestion-dechets-non-pertinent",
    "B7-flux-massique",
    "B7-flux-massique-non-pertinent",
    "C3-reduction-scope1et2",
    "C3-reduction-scope1et2-non-pertinent",
    "C3-reduction-scope3-1",
    "C3-reduction-scope3-2",
    "C3-reduction-scope3-non-applicable",
    "C3-plan-transition-non-pertinent",
    "C3-plan-transition",
    "C3-contribution-plan",
    "C4-aleas",
    "C4-aleas-non-pertinent",
    "C4-impacts-financiers",
    "C4-impacts-financiers-non-applicable",
    "C4-impacts-financiers-non-pertinents",
    "social",
    "B8-effectifs-contrat",
    "B8-effectifs-contrat-partiellement-non-applicable",
    "B8-effectifs-pays",
    "B9-accidents",
    "B10-remuneration",
    "B10-remuneration-partiellement-non-applicable",
    "B10-formation",
    "C5-effectifs",
    "C5-effectifs-non-applicable",
    "C5-effectifs-non-pertinent",
    "C6-droits-homme",
    "C6-droits-homme-non-pertinent",
    "C6-autres-themes",
    "C7-incidents-graves",
    "C7-incidents-graves-non-pertinent",
    "C7-actions",
    "C7-actions-non-pertinent",
    "gouvernance",
    "B11-condamnations",
    "B11-condamnations-non-pertinent",
    "C8-chiffre-affaires",
    "C8-chiffre-affaires-non-pertinent",
    "C8-accord-paris",
    "C9-mixite",
    "C9-mixite-non-pertinent",
    "fin",
]


def pptx_response(presentation, filename):
    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)

    response = HttpResponse(
        buffer.read(),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    response["Content-Disposition"] = f"filename={filename}"

    return response


def find_slide(presentation, name):
    index = find_slide_index(presentation, name)
    return presentation.slides[index]


def find_slide_index(presentation, name):
    return NOMS_DIAPO.index(name)


def find_shape(shapes, name):
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == 6:  # GROUPE
            result = find_shape(shape.shapes, name)
            if result:
                return result


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def remove_slide(presentation, slide_index):
    """supprime la diapo et sa référence pour être plus portable"""
    xml_slides = presentation.slides._sldIdLst
    slide_elem = xml_slides[slide_index]
    rId = slide_elem.get(qn("r:id"))
    presentation.part.drop_rel(rId)
    xml_slides.remove(slide_elem)
