from pathlib import Path

import pytest
from django.conf import settings
from django.urls import reverse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from vsme.export_pptx import export_indicateurs
from vsme.export_pptx import export_sommaire
from vsme.export_pptx import find_shape
from vsme.export_pptx import formate_valeur
from vsme.export_pptx import selectionne_diapos_modules_complets
from vsme.export_pptx import selectionne_diapos_non_applicables
from vsme.export_pptx import selectionne_diapos_non_pertinents
from vsme.models import ExigenceDePublication
from vsme.models import Indicateur
from vsme.models import RapportVSME


def test_telechargement_d_un_rapport_vsme_au_format_pptx_est_privé(
    client, entreprise_factory, alice, bob
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    url = f"/vsme/{rapport_vsme.id}/export/pptx"

    # non connecté
    response = client.get(url)

    assert response.status_code == 302
    connexion_url = reverse("users:login")
    assert response.url == f"{connexion_url}?next={url}"

    # connecté mais non membre de l'entreprise
    client.force_login(bob)

    response = client.get(f"/vsme/{rapport_vsme.id}/export/pptx")

    assert response.status_code == 403


def test_telechargement_d_un_rapport_vsme_au_format_pptx(
    client, entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)

    client.force_login(alice)

    response = client.get(f"/vsme/{rapport_vsme.id}/export/pptx")

    assert response["Content-Disposition"] == "filename=vsme.pptx"
    assert (
        response["content-type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@pytest.mark.parametrize(
    "visibilite_selon_module", [("base", False), ("complet", True)]
)
def test_export_du_sommaire_pptx_selon_le_module_selectionné(
    visibilite_selon_module, entreprise_factory, alice
):
    choix_module, visibilite = visibilite_selon_module
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-a",  # choix du module
        data={"choix_module": choix_module},
    )
    indicateur.save()
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_sommaire(rapport_vsme, presentation)

    shapes = presentation.slides[2].shapes
    assert (
        bool(find_shape(shapes, "Round Same-side Corner of Rectangle 22")) == visibilite
    )
    assert (
        bool(find_shape(shapes, "Round Same-side Corner of Rectangle 23")) == visibilite
    )


def test_export_pptx_d_un_champ_nombre_entier(entreprise_factory, alice):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-e-iii",  # bilan
        data={"bilan": 12345},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[4].shapes
    for shape in shapes:
        if shape.name == "B1-24-e-iii":
            assert shape.text_frame.paragraphs[1].runs[0].text == "12345 euros"


def test_export_pptx_d_un_champ_texte_long(entreprise_factory, alice):
    """il peut avoir des 'paragraphs' vides à la fin qu'il faut ignorer"""
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B7-37",  # Principes de l'économie circulaire
        data={"non_pertinent": False, "economie_circulaire": "PRINCIPES"},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[37].shapes
    for shape in shapes:
        if shape.name == "Rounded Rectangle 8":
            assert shape.text_frame.paragraphs[1].runs[0].text == "PRINCIPES"


def test_export_pptx_d_un_champ_choix_unique(entreprise_factory, alice):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-a",  # Base d’établissement
        data={"choix_module": "complet"},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[4].shapes
    for shape in shapes:
        if shape.name == "B1-24-a":
            assert shape.text_frame.paragraphs[1].runs[0].text == "Module complet"


def test_export_pptx_d_un_champ_choix_binaire(entreprise_factory, alice):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-e-i",  # Forme juridique de l’entreprise
        data={"forme_juridique": "57", "coopérative": True},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[4].shapes
    for shape in shapes:
        if shape.name == "B1-24-e-i_coopérative":
            assert shape.text_frame.paragraphs[1].runs[0].text == "OUI"


def test_export_pptx_d_un_champ_choix_binaire__forme_juridique(
    entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-e-i",  # Forme juridique de l’entreprise
        data={"forme_juridique": "57", "coopérative": True},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[4].shapes
    for shape in shapes:
        if shape.name == "B1-24-e-i_forme_juridique":
            assert (
                shape.text_frame.paragraphs[1].runs[0].text
                == "Société par actions simplifiée (SAS)"
            )


def test_export_pptx_d_un_champ_choix_multiple(entreprise_factory, alice):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-e-ii",  # Code(s) NACE
        data={"nace": ["03.11", "03.21", "03.30"]},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[4].shapes
    for shape in shapes:
        if shape.name == "B1-24-e-ii":
            assert (
                shape.text_frame.paragraphs[1].runs[0].text
                == "Pêche en mer, Aquaculture en mer, Activités de soutien à la pêche et l’aquaculture"
            )


def test_export_pptx_d_un_champ_tableau(entreprise_factory, alice):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-e-vii",  # Liste des sites
        data={
            "sites": [
                {
                    "id_site": 1,
                    "nom_site": "usine",
                    "adresse": "rue du secondaire",
                    "code_postal": "1028",
                    "ville": "Tirana",
                    "pays": "ALB",
                    "geolocalisation": "[1,2]",
                },
                {
                    "id_site": 2,
                    "nom_site": "bureaux",
                    "adresse": "rue du tertiaire",
                    "code_postal": "33000",
                    "ville": "Bordeaux",
                    "pays": "FRA",
                    "geolocalisation": "[3,4]",
                },
            ]
        },
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[6].shapes
    for shape in shapes:
        if shape.name == "B1-25":
            tableau = shape.table
            assert tableau.cell(1, 0).text == "1"
            assert tableau.cell(1, 1).text == "usine"
            assert tableau.cell(1, 2).text == "rue du secondaire"
            assert tableau.cell(1, 3).text == "1028"
            assert tableau.cell(1, 4).text == "Tirana"
            assert tableau.cell(1, 5).text == "ALBANIE"
            assert tableau.cell(1, 6).text == "[1,2]"
            assert tableau.cell(2, 0).text == "2"
            assert tableau.cell(2, 1).text == "bureaux"
            assert tableau.cell(2, 2).text == "rue du tertiaire"
            assert tableau.cell(2, 3).text == "33000"
            assert tableau.cell(2, 4).text == "Bordeaux"
            assert tableau.cell(2, 5).text == "FRANCE"
            assert tableau.cell(2, 6).text == "[3,4]"
            paragraph = tableau.cell(2, 6).text_frame.paragraphs[0]
            assert paragraph.font.size == Pt(10)


def test_export_pptx_d_un_champ_tableau_à_lignes_fixes_avec_choix_binaire_radio(
    entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B2-26",  # Déclaration des pratiques et politiques de durabilité
        data={
            "non_pertinent": False,
            "declaration_durabilite": {
                "changement_climatique": {
                    "pratiques": False,
                    "accessibles": False,
                    "cibles": True,
                },
                "pollution": {"pratiques": True, "accessibles": True, "cibles": False},
                "eau": {"pratiques": False, "accessibles": False, "cibles": True},
                "biodiversite": {
                    "pratiques": True,
                    "accessibles": True,
                    "cibles": False,
                },
                "economie_circulaire": {
                    "pratiques": False,
                    "accessibles": False,
                    "cibles": True,
                },
                "personnel": {"pratiques": True, "accessibles": True, "cibles": False},
                "travailleurs": {
                    "pratiques": False,
                    "accessibles": False,
                    "cibles": True,
                },
                "communautes": {
                    "pratiques": True,
                    "accessibles": True,
                    "cibles": False,
                },
                "consommateurs": {
                    "pratiques": False,
                    "accessibles": False,
                    "cibles": True,
                },
                "conduite_affaires": {
                    "pratiques": True,
                    "accessibles": True,
                    "cibles": False,
                },
            },
        },
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[11].shapes
    for shape in shapes:
        if shape.name == "Table 5":
            tableau = shape.table
            assert tableau.cell(1, 0).text == "Changement climatique"
            assert tableau.cell(1, 1).text == "NON"
            assert tableau.cell(1, 2).text == "NON"
            assert tableau.cell(1, 3).text == "OUI"
            assert tableau.cell(2, 0).text == "Pollution"
            assert tableau.cell(2, 1).text == "OUI"
            assert tableau.cell(2, 2).text == "OUI"
            assert tableau.cell(2, 3).text == "NON"
            para_oui = tableau.cell(1, 3).text_frame.paragraphs[0]
            para_non = tableau.cell(1, 1).text_frame.paragraphs[0]
            assert para_oui.font.color.rgb == RGBColor(0x14, 0x5C, 0x30)
            assert para_non.font.color.rgb == RGBColor(0xC0, 0x00, 0x00)
            assert para_oui.alignment == para_non.alignment == PP_ALIGN.CENTER
            assert para_oui.font.size == para_non.font.size == Pt(10)


def test_export_pptx_d_un_champ_tableau_à_lignes_fixes_avec_données_vides(
    entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B3-30-p2",  # Déclaration des pratiques et politiques de durabilité
        data={
            "non_pertinent": False,
            "estimation_emissions_GES_scope_3": {
                "biens_et_services_achetes": {"emissions_brutes_GES": "23"},
                "biens_investissement": {"emissions_brutes_GES": "44"},
                "activites_secteurs_combustibles": {"emissions_brutes_GES": None},
                "transport_distribution_amont": {"emissions_brutes_GES": None},
                "dechets_produits": {"emissions_brutes_GES": None},
                "voyages_affaires": {"emissions_brutes_GES": None},
                "deplacements_domicile_travail": {"emissions_brutes_GES": None},
                "actifs_loues_amont": {"emissions_brutes_GES": None},
                "acheminement_aval": {"emissions_brutes_GES": None},
                "transformation_produits_vendus": {"emissions_brutes_GES": None},
                "utilisation_produits_vendus": {"emissions_brutes_GES": None},
                "traitement_fin_de_vie_produits_vendus": {"emissions_brutes_GES": None},
                "actifs_loues_aval": {"emissions_brutes_GES": None},
                "franchises": {"emissions_brutes_GES": None},
                "investissements": {"emissions_brutes_GES": None},
            },
            "total_estimation_emissions_GES_scope_3": {
                "total_scope_3": {"total_emissions_brutes_GES": None},
                "total_scopes_1_2_3": {"total_emissions_brutes_GES": None},
            },
        },
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[24].shapes
    for shape in shapes:
        if shape.name == "Table 18":
            tableau = shape.table
            assert tableau.cell(1, 0).text == "1. Biens et services achetés"
            assert tableau.cell(1, 1).text == "23"
            assert (
                tableau.cell(3, 0).text
                == "3. Activités relevant des secteurs des combustibles et de l'énergie"
            )
            assert tableau.cell(3, 1).text == "0"


def test_export_pptx_d_un_champ_tableau_à_lignes_fixes_sur_plusieurs_diapos(
    entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="C2-48",
        data={
            "description_durabilite": {
                "changement_climatique": {
                    "description_pratiques": "a",
                    "description_cibles": "b",
                    "niveau_hierarchique": "c",
                },
                "pollution": {
                    "description_pratiques": "d",
                    "description_cibles": "e",
                    "niveau_hierarchique": "f",
                },
                "eau": {
                    "description_pratiques": "g",
                    "description_cibles": "h",
                    "niveau_hierarchique": "i",
                },
                "biodiversite": {
                    "description_pratiques": "j",
                    "description_cibles": "k",
                    "niveau_hierarchique": "l",
                },
                "economie_circulaire": {
                    "description_pratiques": "m",
                    "description_cibles": "n",
                    "niveau_hierarchique": "o",
                },
                "personnel": {
                    "description_pratiques": "p",
                    "description_cibles": "q",
                    "niveau_hierarchique": "r",
                },
                "travailleurs": {
                    "description_pratiques": "s",
                    "description_cibles": "t",
                    "niveau_hierarchique": "u",
                },
                "communautes": {
                    "description_pratiques": "v",
                    "description_cibles": "w",
                    "niveau_hierarchique": "x",
                },
                "consommateurs": {
                    "description_pratiques": "y",
                    "description_cibles": "z",
                    "niveau_hierarchique": "aa",
                },
                "conduite_affaires": {
                    "description_pratiques": "bb",
                    "description_cibles": "cc",
                    "niveau_hierarchique": "dd",
                },
            }
        },
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[16].shapes
    for shape in shapes:
        if shape.name == "Table 8":
            tableau = shape.table
            assert tableau.cell(1, 0).text == "Changement climatique"
            assert tableau.cell(1, 1).text == "a"
            assert tableau.cell(2, 0).text == "Pollution"
            assert tableau.cell(2, 1).text == "d"
    shapes = presentation.slides[17].shapes
    for shape in shapes:
        if shape.name == "Table 8":
            tableau = shape.table
            assert tableau.cell(1, 0).text == "Eau et ressources aquatiques et marines"
            assert tableau.cell(1, 1).text == "g"
            assert tableau.cell(2, 0).text == "Biodiversité et écosystèmes"
            assert tableau.cell(2, 1).text == "j"
    shapes = presentation.slides[20].shapes
    for shape in shapes:
        if shape.name == "Table 8":
            tableau = shape.table
            assert tableau.cell(1, 0).text == "Consommateurs et utilisateurs finaux"
            assert tableau.cell(1, 1).text == "y"
            assert (
                tableau.cell(2, 0).text
                == "Conduite des affaires et lutte contre la corruption"
            )
            assert tableau.cell(2, 1).text == "bb"


def test_selectionne_diapos_non_applicables_d_un_indicateur_non_applicable(
    entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    Indicateur.objects.create(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-e-i",  # Forme juridique de l’entreprise
        data={"forme_juridique": "57", "coopérative": False},
    )

    diapos_a_supprimer = selectionne_diapos_non_applicables(rapport_vsme)

    # "B2-26-p1" "Participation effective à la gouvernance" est non applicable
    # donc suppression de ‘diapo’
    schema_indicateur = list(
        ExigenceDePublication.par_indicateur_schema_id("B2-26-p1")
        .load_json_schema()
        .values()
    )[0]
    index_diapo_a_supprimer = schema_indicateur["champs"][0]["export_pptx"][
        "diapo"
    ]  # 10
    assert index_diapo_a_supprimer in diapos_a_supprimer


def test_selectionne_diapos_non_applicables_d_un_indicateur_applicable(
    entreprise_factory, alice
):
    # il faut supprimer la diapo non applicable
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    Indicateur.objects.create(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-e-i",  # Forme juridique de l’entreprise
        data={"forme_juridique": "57", "coopérative": True},
    )

    diapos_a_supprimer = selectionne_diapos_non_applicables(rapport_vsme)

    # "B2-26-p1" "Participation effective à la gouvernance" est applicable
    # donc suppression de 'diapo_non_applicable'
    schema_indicateur = list(
        ExigenceDePublication.par_indicateur_schema_id("B2-26-p1")
        .load_json_schema()
        .values()
    )[0]
    index_diapo_a_supprimer = schema_indicateur["champs"][0]["export_pptx"][
        "diapo_non_applicable"
    ]  # 9
    assert index_diapo_a_supprimer in diapos_a_supprimer


@pytest.mark.parametrize(
    "valeur, champ, attendu",
    [
        (42, {"type": "nombre_entier"}, "42"),
        (None, {"type": "nombre_entier"}, "0"),
        (0, {"type": "nombre_entier"}, "0"),
        (42, {"type": "nombre_entier", "unité": "kg"}, "42 kg"),
        (None, {"type": "nombre_entier", "unité": "kg"}, "0 kg"),
        (0, {"type": "nombre_entier", "unité": "kg"}, "0 kg"),
        (3.14, {"type": "nombre_decimal"}, "3.14"),
        (None, {"type": "nombre_decimal"}, "0"),
        (3.14, {"type": "nombre_decimal", "unité": "tonnes"}, "3.14 tonnes"),
        (None, {"type": "nombre_decimal", "unité": "tonnes"}, "0 tonnes"),
        (1, {"type": "auto_id"}, "1"),
        (
            "-50.0%",
            {"type": "nombre_decimal", "unité": "%", "calculé": True},
            "-50.0 %",
        ),
        (None, {"type": "auto_id"}, "0"),
        (True, {"type": "choix_binaire"}, "OUI"),
        (False, {"type": "choix_binaire"}, "NON"),
    ],
)
def test_formate_valeur(valeur, champ, attendu):
    assert formate_valeur(valeur, champ) == attendu


def test_export_pptx_d_un_indicateur_non_pertinent(entreprise_factory, alice):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B7-37",
        data={"non_pertinent": True, "economie_circulaire": "PRINCIPES"},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[37].shapes
    for shape in shapes:
        if shape.name == "Rounded Rectangle 8":
            assert shape.text_frame.paragraphs[1].runs[0].text != "PRINCIPES"


def test_export_pptx_d_un_indicateur_non_pertinent_inscrit_le_texte_non_pertinent(
    entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B2-26-p1",
        data={
            "non_pertinent": True,
            "participation_gouvernance": "texte de participation",
        },
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[9].shapes  # diapo 10
    for shape in shapes:
        if shape.name == "Rounded Rectangle 1":
            paragraphe = None
            for p in shape.text_frame.paragraphs:
                if p.runs:
                    paragraphe = p
            assert paragraphe.runs[0].text == "Non pertinent"


def test_export_pptx_d_un_indicateur_non_pertinent_inscrit_le_texte_si_pertinent(
    entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B6-36",
        data={"non_pertinent": True, "consommation_eau": "500"},
    )
    chemin_pptx = Path(settings.BASE_DIR, "vsme/exports/vsme.pptx")
    presentation = Presentation(chemin_pptx)

    export_indicateurs([indicateur], presentation)

    shapes = presentation.slides[35].shapes  # diapo 36
    for shape in shapes:
        if shape.name == "Rounded Rectangle 15":
            paragraphe = None
            for p in shape.text_frame.paragraphs:
                if p.runs:
                    paragraphe = p
            assert paragraphe.runs[0].text == (
                "Les processus de production de l'entreprise n'entraînent pas une consommation d'eau importante"
            )


@pytest.mark.parametrize(
    "choix_module, diapos_modules_complets_vides",
    [("complet", True), ("base", False)],
)
def test_selectionne_diapos_modules_complets(
    choix_module, diapos_modules_complets_vides, entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    Indicateur.objects.create(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-a",
        data={"choix_module": choix_module},
    )

    diapos = selectionne_diapos_modules_complets(rapport_vsme)

    if diapos_modules_complets_vides:
        assert diapos == set()
    else:
        # diapos C1 (diapo simple) et C2 (multidiapos)
        assert {14, 15, 17, 18, 19, 20, 21}.issubset(diapos)


@pytest.mark.parametrize(
    "non_pertinent, diapo_a_supprimer",
    [(True, {12}), (False, {13})],
)
def test_selectionne_diapos_non_pertinents(
    non_pertinent, diapo_a_supprimer, entreprise_factory, alice
):
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    indicateur = Indicateur(
        rapport_vsme=rapport_vsme,
        schema_id="B2-26",
        data={"non_pertinent": non_pertinent},
    )

    assert selectionne_diapos_non_pertinents([indicateur]) == diapo_a_supprimer


def test_selectionne_diapos_non_applicables_C4_57_et_C4_58_non_applicables(
    entreprise_factory, alice
):
    # Module de base → C4-57 et C4-58 non applicables.
    # C4-57 : diapo_non_pertinent (54) == diapo_non_applicable (54) → seule diapo (53) supprimée.
    # C4-58 : diapo_non_pertinent (57) != diapo_non_applicable (56) → diapo (55) et diapo_non_pertinent (57) supprimées.
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    Indicateur.objects.create(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-a",
        data={"choix_module": "base"},
    )

    diapos_a_supprimer = selectionne_diapos_non_applicables(rapport_vsme)

    assert 53 in diapos_a_supprimer
    assert 54 not in diapos_a_supprimer
    assert 55 in diapos_a_supprimer
    assert 56 not in diapos_a_supprimer
    assert 57 in diapos_a_supprimer


def test_selectionne_diapos_non_applicables_C4_57_et_C4_58_applicables(
    entreprise_factory, alice
):
    # Module complet → C4-57 et C4-58 applicables.
    # C4-58 nécessite des risques climatiques dans C4-57 pour être applicable.
    # Seule diapo_non_applicable est supprimée (54 et 56).
    entreprise = entreprise_factory(utilisateur=alice)
    rapport_vsme = RapportVSME.objects.create(entreprise=entreprise, annee=2026)
    Indicateur.objects.create(
        rapport_vsme=rapport_vsme,
        schema_id="B1-24-a",
        data={"choix_module": "complet"},
    )
    Indicateur.objects.create(
        rapport_vsme=rapport_vsme,
        schema_id="C4-57",
        data={
            "aleas_et_risques_climatiques": [
                {"id_risque": 1, "description": "inondation"}
            ]
        },
    )

    diapos_a_supprimer = selectionne_diapos_non_applicables(rapport_vsme)

    assert 53 not in diapos_a_supprimer
    assert 54 in diapos_a_supprimer
    assert 55 not in diapos_a_supprimer
    assert 56 in diapos_a_supprimer
    assert 57 not in diapos_a_supprimer
