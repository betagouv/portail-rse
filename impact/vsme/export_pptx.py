import copy

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from utils.pays import CODES_PAYS_ISO_3166_1
from utils.pptx import find_shape
from utils.pptx import remove_shape
from utils.pptx import remove_slide
from vsme.export_xlsx import formate_valeur as formate_valeur_xlsx
from vsme.forms import THEMATIQUES_DURABILITE
from vsme.models import ExigenceDePublication
from vsme.models import EXIGENCES_DE_PUBLICATION


def export_rapport_vsme(rapport_vsme, presentation):
    indicateurs = rapport_vsme.indicateurs.all()
    export_couverture(rapport_vsme, presentation)
    export_sommaire(rapport_vsme, presentation)
    export_indicateurs(indicateurs, presentation)
    supprime_diapos_inutiles(indicateurs, rapport_vsme, presentation)


def export_couverture(rapport_vsme, presentation):
    shapes = presentation.slides[0].shapes
    shape_titre = find_shape(shapes, "ZoneTexte 9")
    shape_titre.text_frame.paragraphs[0].runs[
        0
    ].text = rapport_vsme.entreprise.denomination
    shape_titre.text_frame.paragraphs[1].runs[
        0
    ].text = f"Rapport VSME {rapport_vsme.annee}"


def export_sommaire(rapport_vsme, presentation):
    if rapport_vsme.choix_module == "base":
        NUMERO_DIAPO_SOMMAIRE = 3
        _remove_shapes_from_slide(
            presentation,
            NUMERO_DIAPO_SOMMAIRE,
            ("22", "23", "39", "40", "47", "48", "49", "54", "55", "56"),
        )

        NUMERO_DIAPO_INFOS = 4
        _remove_shapes_from_slide(presentation, NUMERO_DIAPO_INFOS, ("22", "23"))

        NUMERO_DIAPO_ENV = 23
        _remove_shapes_from_slide(presentation, NUMERO_DIAPO_ENV, ("31", "32"))

        NUMERO_DIAPO_SOCIAL = 58
        _remove_shapes_from_slide(presentation, NUMERO_DIAPO_SOCIAL, ("10", "11", "12"))

        NUMERO_DIAPO_GOUV = 75
        _remove_shapes_from_slide(presentation, NUMERO_DIAPO_GOUV, ("17", "18", "20"))


def _remove_shapes_from_slide(presentation, slide_num, shape_nums):
    shapes = presentation.slides[slide_num - 1].shapes
    for num_shape in shape_nums:
        shape = find_shape(shapes, f"Round Same-side Corner of Rectangle {num_shape}")
        remove_shape(shape)


def supprime_diapos_inutiles(indicateurs, rapport_vsme, presentation):
    diapos_a_supprimer = (
        selectionne_diapos_non_pertinents(indicateurs)
        | selectionne_diapos_non_applicables(rapport_vsme)
        | selectionne_diapos_modules_complets(rapport_vsme)
    )
    for diapo_a_supprimer in sorted(diapos_a_supprimer, reverse=True):
        index_diapo = diapo_a_supprimer - 1
        remove_slide(presentation, index_diapo)


def selectionne_diapos_non_pertinents(indicateurs):
    diapos_a_supprimer = set()
    for indicateur in indicateurs:
        for champ in indicateur.schema["champs"]:
            if "export_pptx" not in champ:
                continue
            export_pptx = champ["export_pptx"]
            if "diapo_non_pertinent" not in export_pptx:
                continue
            if indicateur.est_non_pertinent:
                diapos_a_supprimer.add(export_pptx["diapo"])
            else:
                diapos_a_supprimer.add(export_pptx["diapo_non_pertinent"])
    return diapos_a_supprimer


def selectionne_diapos_non_applicables(rapport_vsme):
    diapos_a_supprimer = set()

    tous_les_indicateur_schema_ids = [
        schema_id
        for exigence in EXIGENCES_DE_PUBLICATION.values()
        for schema_id in exigence.indicateurs()
    ]

    for indicateur_schema_id in tous_les_indicateur_schema_ids:
        schema_exigence = ExigenceDePublication.par_indicateur_schema_id(
            indicateur_schema_id
        ).load_json_schema()
        schema_indicateur = schema_exigence[indicateur_schema_id]
        if rapport_vsme.indicateur_est_applicable(indicateur_schema_id)[0]:
            # supprimer les diapos non applicables
            for champ in schema_indicateur["champs"]:
                if (
                    "export_pptx" in champ
                    and "diapo_non_applicable" in champ["export_pptx"]
                ):
                    diapo_a_supprimer = champ["export_pptx"]["diapo_non_applicable"]
                    diapos_a_supprimer.add(diapo_a_supprimer)
        else:
            # supprimer les diapos applicables
            for champ in schema_indicateur["champs"]:
                if (
                    "export_pptx" in champ
                    and "diapo_non_applicable" in champ["export_pptx"]
                ):
                    export_pptx = champ["export_pptx"]
                    if "diapo" in export_pptx:
                        diapos_a_supprimer.add(export_pptx["diapo"])
                    for multidiapo in export_pptx.get("multidiapos", []):
                        diapos_a_supprimer.add(multidiapo["diapo"])

    return diapos_a_supprimer


def selectionne_diapos_modules_complets(rapport_vsme):
    if rapport_vsme.choix_module == "complet":
        return set()
    diapos_a_supprimer = set()
    for code, exigence in EXIGENCES_DE_PUBLICATION.items():
        if not code.startswith("C"):
            continue
        schema = exigence.load_json_schema()
        for schema_indicateur in schema.values():
            for champ in schema_indicateur["champs"]:
                export_pptx = champ.get("export_pptx")
                if not export_pptx:
                    continue
                if "diapo" in export_pptx:
                    diapos_a_supprimer.add(export_pptx["diapo"])
                if "diapo_non_pertinent" in export_pptx:
                    diapos_a_supprimer.add(export_pptx["diapo_non_pertinent"])
                for multidiapo in export_pptx.get("multidiapos", []):
                    diapos_a_supprimer.add(multidiapo["diapo"])
    return diapos_a_supprimer


def export_indicateurs(indicateurs, presentation):
    for indicateur in indicateurs:
        _export_indicateur(indicateur, presentation)


def _export_indicateur(indicateur, presentation):
    if indicateur.est_non_pertinent:
        return

    for champ in indicateur.schema["champs"]:
        if "export_pptx" not in champ:
            continue
        data = indicateur.data.get(champ["id"])
        export_pptx = champ["export_pptx"]
        if "multidiapos" in export_pptx:
            _export_champ_multidiapos(champ, data, export_pptx, presentation)
        else:
            _export_champ_monodiapo(champ, data, export_pptx, presentation)


def _export_champ_multidiapos(champ, data, export_pptx, presentation):
    multidiapos = export_pptx["multidiapos"]
    lignes = champ["lignes"]
    if lignes == "THEMATIQUES_DURABILITE":
        label_to_id = {v: k for k, v in THEMATIQUES_DURABILITE.items()}
    else:
        label_to_id = {ligne["label"]: ligne["id"] for ligne in lignes}
    colonnes = champ["colonnes"]
    for diapo_info in multidiapos:
        diapo = presentation.slides[diapo_info["diapo"] - 1]
        shape = find_shape(diapo.shapes, diapo_info["shape"])
        if not shape:
            continue
        table = shape.table
        for row_index in range(1, len(table.rows)):
            label = table.cell(row_index, 0).text.strip()
            if label not in label_to_id:
                continue
            ligne_id = label_to_id[label]
            data_ligne = data[ligne_id]
            for col_index, colonne in enumerate(colonnes):
                data_cellule = data_ligne.get(colonne["id"])
                cell = table.cell(row_index, col_index + 1)
                cell.text = formate_valeur(data_cellule, colonne)
                _appliquer_style_cellule(cell, data_cellule, colonne)


def _export_champ_monodiapo(champ, data, export_pptx, presentation):
    numero_diapo = export_pptx["diapo"]
    diapo = presentation.slides[numero_diapo - 1]
    nom_shape = export_pptx["shape"]
    shape = find_shape(diapo.shapes, nom_shape)
    if not shape:
        print("Non trouvé", champ["id"], nom_shape)
    _export_champ(champ, data, shape)


def _export_champ(champ, data, shape):
    type_indicateur = champ["type"]
    match type_indicateur:
        case "choix_multiple":
            _export_choix_multiple(champ, data, shape)
        case "tableau":
            _export_tableau(champ, data, shape)
        case "tableau_lignes_fixes":
            _export_tableau_lignes_fixes(champ, data, shape)
        case _:
            _export_simple(champ, data, shape)


def _export_simple(champ, data, shape):
    valeur = formate_valeur(data, champ)
    paragraphs = shape.text_frame.paragraphs
    paragraphe_a_remplir = None
    for paragraphe in shape.text_frame.paragraphs:
        if paragraphe.runs:
            paragraphe_a_remplir = paragraphe
    paragraphe_a_remplir.runs[0].text = str(valeur)


def _export_choix_multiple(champ, data, shape):
    valeurs = (formate_valeur(valeur, champ) for valeur in data)
    shape.text_frame.paragraphs[1].runs[0].text = ", ".join(valeurs)


def formate_valeur(valeur, champ):
    match champ["type"]:
        case "nombre_entier" | "nombre_decimal" | "auto_id":
            if isinstance(valeur, str) and valeur.endswith("%"):
                valeur = valeur[:-1]
            if "unité" in champ:
                unite = champ["unité"]
                valeur_formatee = f"{valeur} {unite}" if valeur else f"0 {unite}"
            else:
                valeur_formatee = f"{valeur}" if valeur else "0"
            return valeur_formatee
        case _:
            return str(formate_valeur_xlsx(valeur, champ))


def _export_tableau(champ, data, shape):
    table = shape.table
    _formate_hauteur_tableau(data, table)

    for index_ligne, ligne in enumerate(data, start=1):
        for index_data, data_cellule in enumerate(ligne.values()):
            schema_colonne = champ["colonnes"][index_data]
            cell = table.cell(index_ligne, index_data)
            cell.text = formate_valeur(data_cellule, schema_colonne)
            _appliquer_style_cellule(cell, data_cellule, schema_colonne)


def _formate_hauteur_tableau(data, table):
    nombre_lignes_a_garder = 2  # titre et première ligne de données
    trs_a_supprimer = [
        table.rows[i]._tr for i in range(nombre_lignes_a_garder, len(table.rows))
    ]
    for tr in trs_a_supprimer:
        table._tbl.remove(tr)

    nombre_lignes_a_ajouter = len(data) - 1  # première ligne déjà présente donc -1
    for _ in range(nombre_lignes_a_ajouter):
        _ajoute_ligne_tableau(table)


def _ajoute_ligne_tableau(table):
    tr_source = table.rows[1]._tr
    tr_nouveau = copy.deepcopy(tr_source)
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for t in tr_nouveau.findall(f".//{{{ns}}}t"):
        t.text = ""
    table._tbl.append(tr_nouveau)


def _export_tableau_lignes_fixes(champ, data, shape):
    lignes = champ["lignes"]
    colonnes = champ["colonnes"]
    match lignes:
        case "PAYS":
            table = shape.table
            _formate_hauteur_tableau(data, table)
            for offset_ligne, (code_pays, data_ligne) in enumerate(data.items()):
                cell = table.cell(offset_ligne + 1, 0)
                cell.text = CODES_PAYS_ISO_3166_1[code_pays]
                for offset_colonne, colonne in enumerate(colonnes):
                    data_cellule = data_ligne.get(colonne["id"])
                    cell = table.cell(offset_ligne + 1, offset_colonne + 1)
                    cell.text = formate_valeur(data_cellule, colonne)
                    _appliquer_style_cellule(cell, data_cellule, colonne)
        case "RISQUES_CLIMATIQUES":
            table = shape.table
            _formate_hauteur_tableau(data, table)
            for offset_ligne, (id_risque, data_ligne) in enumerate(data.items()):
                cell = table.cell(offset_ligne + 1, 0)
                cell.text = formate_valeur(id_risque, {"type": "auto_id"})
                _appliquer_style_cellule(cell, id_risque, {"type": "auto_id"})

                for offset_colonne, colonne in enumerate(colonnes):
                    data_cellule = data_ligne.get(colonne["id"])
                    cell = table.cell(offset_ligne + 1, offset_colonne + 1)
                    cell.text = formate_valeur(data_cellule, colonne)
                    _appliquer_style_cellule(cell, data_cellule, colonne)
        case "THEMATIQUES_DURABILITE" | list():
            colonnes_ids = [colonne["id"] for colonne in colonnes]
            if lignes == "THEMATIQUES_DURABILITE":
                lignes_ids = list(THEMATIQUES_DURABILITE.keys())
            else:
                lignes_ids = [ligne["id"] for ligne in lignes]

            for id_ligne, data_ligne in data.items():
                offset_ligne = lignes_ids.index(id_ligne)
                for id_colonne, data_cellule in data_ligne.items():
                    offset_colonne = colonnes_ids.index(id_colonne)
                    champ = colonnes[offset_colonne]
                    index_ligne = offset_ligne + 1
                    index_colonne = offset_colonne + 1
                    cell = shape.table.cell(index_ligne, index_colonne)
                    cell.text = formate_valeur(data_cellule, champ)
                    _appliquer_style_cellule(cell, data_cellule, champ)


def _appliquer_style_cellule(cell, data_cellule, champ):
    type_indicateur = champ["type"]
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in cell.text_frame.paragraphs:
        para.font.size = Pt(10)
        para.alignment = PP_ALIGN.CENTER
        if type_indicateur == "choix_binaire_radio":
            ROUGE = RGBColor(192, 0, 0)
            VERT = RGBColor(20, 92, 48)
            if data_cellule:
                para.font.color.rgb = VERT
            else:
                para.font.color.rgb = ROUGE
